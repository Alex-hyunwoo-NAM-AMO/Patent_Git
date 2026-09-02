import argparse
import json
import os
import sys

from pptx import Presentation


def extract_pptx(path):
    prs = Presentation(path)
    slides = []
    media_count = 0
    for idx, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            elif shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
            elif shape.shape_type == 13:
                media_count += 1
        slides.append({"index": idx, "text": "\n".join(parts)})
    return slides, media_count


def extract_pdf(path):
    import pdfplumber

    slides = []
    with pdfplumber.open(path) as pdf:
        for idx, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            slides.append({"index": idx, "text": text.strip()})
    return slides, 0


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("path")
    args = ap.parse_args()

    ext = os.path.splitext(args.path)[1].lower()
    if ext == ".pptx":
        slides, media_count = extract_pptx(args.path)
    elif ext == ".pdf":
        slides, media_count = extract_pdf(args.path)
    else:
        json.dump({"error": f"지원하지 않는 형식: {ext}"}, sys.stdout)
        sys.exit(2)

    full_text = "\n\n".join(s["text"] for s in slides if s["text"])
    json.dump(
        {
            "slides": slides,
            "fullText": full_text,
            "sourcePath": args.path,
            "mediaCount": media_count,
        },
        sys.stdout,
        ensure_ascii=False,
    )


if __name__ == "__main__":
    main()
