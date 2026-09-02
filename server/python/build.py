import argparse
import copy
import json
import os
import sys
from io import BytesIO

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn


def set_cell(table, row, col, value):
    if value is None or value == "":
        return
    cell = table.cell(row, col)
    tf = cell.text_frame
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = str(value)
        for extra in first_para.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        first_para.add_run().text = str(value)
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)


def get_tables(slide):
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            tables.append(shape.table)
    return tables


def fill_basic_and_inventors(slide, sheet, inventors):
    tables = get_tables(slide)
    if len(tables) < 3:
        return
    basic, inv_table, legal = tables[0], tables[1], tables[2]

    b = sheet.get("basic", {})
    set_cell(basic, 0, 1, b.get("title"))
    set_cell(basic, 1, 1, b.get("corporation"))
    set_cell(basic, 1, 3, b.get("completionDate"))
    set_cell(basic, 3, 1, b.get("applications"))

    for i, inv in enumerate(inventors[:4]):
        r = i + 1
        set_cell(inv_table, r, 1, inv.get("name"))
        set_cell(inv_table, r, 2, inv.get("employeeId"))
        set_cell(inv_table, r, 3, inv.get("department"))
        set_cell(inv_table, r, 4, inv.get("position"))
        set_cell(inv_table, r, 5, inv.get("email"))
        set_cell(inv_table, r, 6, inv.get("contributionPct"))

    lg = sheet.get("legal", {})
    set_cell(legal, 0, 1, lg.get("nationalRnd"))
    set_cell(legal, 1, 1, lg.get("externalCoDev"))
    set_cell(legal, 2, 1, lg.get("disclosure"))


def fill_gist(slide, sheet):
    tables = get_tables(slide)
    if not tables:
        return
    g = sheet.get("gist", {})
    gist_table = tables[0]
    set_cell(gist_table, 0, 1, g.get("techField"))
    set_cell(gist_table, 1, 1, g.get("priorArtProblem"))
    set_cell(gist_table, 2, 1, g.get("problemToSolve"))
    set_cell(gist_table, 3, 1, g.get("coreComposition"))
    set_cell(gist_table, 4, 1, g.get("effect"))


def remove_slide(prs, index):
    slides = prs.slides._sldIdLst
    slide_ids = list(slides)
    if index >= len(slide_ids):
        return
    sldId = slide_ids[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    slides.remove(sldId)


_REL_ID_ATTRS = (qn("r:embed"), qn("r:link"), qn("r:id"))


def import_shape_relationship(src_slide, dst_slide, old_rid, rid_map):
    if old_rid in rid_map:
        return rid_map[old_rid]
    src_rel = src_slide.part.rels[old_rid]
    if not src_rel.is_external and src_rel.reltype == RT.IMAGE:
        blob = src_rel.target_part.blob
        _, new_rid = dst_slide.part.get_or_add_image_part(BytesIO(blob))
        rid_map[old_rid] = new_rid
        return new_rid
    if src_rel.is_external:
        new_rid = dst_slide.part.relate_to(src_rel.target_ref, src_rel.reltype, is_external=True)
        rid_map[old_rid] = new_rid
        return new_rid
    new_rid = import_media_blob(dst_slide, src_rel)
    if new_rid is not None:
        rid_map[old_rid] = new_rid
        return new_rid
    raise NotImplementedError(f"지원하지 않는 내부 관계: type={src_rel.reltype!r}")


def import_media_blob(dst_slide, src_rel):
    target = src_rel.target_part
    blob = getattr(target, "blob", None)
    if blob is None:
        return None
    partname = str(getattr(target, "partname", ""))
    ext = os.path.splitext(partname)[1].lstrip(".") or "bin"
    content_type = getattr(target, "content_type", "application/octet-stream")
    pkg = dst_slide.part.package
    new_partname = pkg.next_partname(f"/ppt/media/importedMedia%d.{ext}")
    from pptx.opc.package import Part

    new_part = Part(new_partname, content_type, pkg, blob)
    return dst_slide.part.relate_to(new_part, src_rel.reltype)


def rewrite_relationship_ids(src_slide, dst_slide, cloned_elm, rid_map):
    for elm in cloned_elm.iter():
        for attr in _REL_ID_ATTRS:
            old_rid = elm.get(attr)
            if old_rid is None:
                continue
            elm.set(attr, import_shape_relationship(src_slide, dst_slide, old_rid, rid_map))


def blank_layout(dest_prs):
    for layout in dest_prs.slide_layouts:
        if len(layout.placeholders) == 0:
            return layout
    if len(dest_prs.slide_layouts) > 6:
        return dest_prs.slide_layouts[6]
    return dest_prs.slide_layouts[-1]


def copy_slide(dest_prs, src_slide):
    new_slide = dest_prs.slides.add_slide(blank_layout(dest_prs))
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    rid_map = {}
    for src_shape in src_slide.shapes:
        cloned = copy.deepcopy(src_shape.element)
        rewrite_relationship_ids(src_slide, new_slide, cloned, rid_map)
        new_slide.shapes._spTree.insert_element_before(cloned, "p:extLst")
    return new_slide


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", required=True)
    ap.add_argument("--payload", required=True)
    ap.add_argument("--original", required=True)
    ap.add_argument("--out", required=True)
    args = ap.parse_args()

    with open(args.payload, "r", encoding="utf-8") as f:
        payload = json.load(f)
    sheet = payload.get("sheet", {})
    inventors = payload.get("inventors", [])

    prs = Presentation(args.template)

    slides = list(prs.slides)
    if len(slides) >= 3:
        fill_basic_and_inventors(slides[1], sheet, inventors)
        fill_gist(slides[2], sheet)

    remove_slide(prs, 3)

    if os.path.exists(args.original) and args.original.lower().endswith(".pptx"):
        src_prs = Presentation(args.original)
        for src_slide in src_prs.slides:
            copy_slide(prs, src_slide)

    os.makedirs(os.path.dirname(os.path.abspath(args.out)), exist_ok=True)
    prs.save(args.out)
    json.dump({"ok": True, "out": args.out, "slideCount": len(list(prs.slides))}, sys.stdout, ensure_ascii=False)


if __name__ == "__main__":
    main()
